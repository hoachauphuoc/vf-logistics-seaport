"""
VF Logistics - AI-Powered Enterprise Seaport Platform
Professional PowerPoint Generator
Snowflake CoCo CLI Hackathon 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x0F, 0x11, 0x17)
CARD_BG = RGBColor(0x1A, 0x1F, 0x2E)
ACCENT_BLUE = RGBColor(0x29, 0xB5, 0xE8)
ACCENT_GREEN = RGBColor(0x22, 0xC5, 0x5E)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_YELLOW = RGBColor(0xEA, 0xB3, 0x08)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, left=Inches(0.8), top=Inches(0.5), width=Inches(8.5), font_size=Pt(32), color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = color
    return txBox

def add_text(slide, text, left=Inches(0.8), top=Inches(1.5), width=Inches(8.5), height=Inches(5), font_size=Pt(16), color=LIGHT_GRAY, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)
    return txBox

def add_bullet(slide, items, left=Inches(0.8), top=Inches(1.8), width=Inches(8.5), font_size=Pt(15), color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_kpi_box(slide, label, value, left, top, width=Inches(2.2), color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = MID_GRAY
    p2.alignment = PP_ALIGN.CENTER

def add_table(slide, headers, rows, left=Inches(0.8), top=Inches(2.0), col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [Inches(8.5 / n_cols)] * n_cols
    
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, sum(col_widths), Inches(0.4 * n_rows))
    table = table_shape.table
    
    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    
    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = LIGHT_GRAY

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: Title ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)
    add_title(slide, "VF Logistics", top=Inches(1.8), font_size=Pt(48), color=ACCENT_BLUE)
    add_text(slide, "Document to SAP in 10 Seconds", top=Inches(2.8), font_size=Pt(28), color=WHITE)
    add_text(slide, "Track 1: Workflow Automation\nFrom document chaos to fully automated ERP posting\n\nTeam SORA | Built 100% with Snowflake CoCo CLI\nSnowflake CoCo CLI Hackathon 2026", top=Inches(4.0), font_size=Pt(14), color=MID_GRAY)
    
    # ========== SLIDE 2: Problem ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "The Problem: Maritime Document Chaos")
    add_kpi_box(slide, "Documents per vessel", "200+", Inches(0.8), Inches(1.5))
    add_kpi_box(slide, "Manual classification", "15-30 min", Inches(3.3), Inches(1.5))
    add_kpi_box(slide, "Compliance error cost", "$5K-$50K", Inches(5.8), Inches(1.5), color=ACCENT_RED)
    add_bullet(slide, [
        "Bills of Lading, Invoices, Certificates, EDI — all processed manually",
        "No automated cross-checking between related documents",
        "Fraud & duplicates undetected until cargo mismatch at destination",
        "Compliance errors cause expensive port delays and fines",
    ], top=Inches(3.2))
    
    # ========== SLIDE 3: Solution ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Solution: Automated Workflow Pipeline")
    add_text(slide, "Upload > Classify > Cross-Check > Compliance > Fraud Scan > SAP Post\n(Mendix)   (Cortex AI)  (Rule+AI)      (SQL Rules)    (Pattern)     (Auto)", top=Inches(1.5), font_size=Pt(13), color=ACCENT_BLUE)
    add_text(slide, "KEY: IT TAKES ACTION", top=Inches(2.8), font_size=Pt(18), color=ACCENT_GREEN, bold=True)
    add_bullet(slide, [
        "Not a dashboard. It classifies, validates, blocks, and posts — automatically",
        "Zero-hallucination: deterministic procedures for financial data",
        "Hybrid AI: rules first (free), AI only for edge cases (90% savings)",
        "End-to-end: document to SAP posting, fully automated",
        "Ultra-low cost: ~$0.001 per document (llama3-8b)",
    ], top=Inches(3.5))
    
    # ========== SLIDE 4: Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Technical Architecture")
    add_text(slide, 
        "MENDIX (Low-Code UI)          SNOWFLAKE (AI + Data Platform)\n"
        "--------------------          ----------------------------\n"
        "  Document Upload               Cortex AI Functions\n"
        "  Workflow Engine                16 Stored Procedures\n"
        "  Async Task Queue              228+ Reference Records\n"
        "  Alert Dashboard               Streamlit Dashboard\n"
        "  SAP Posting UI                Marketplace Weather\n"
        "                                Snowpark Pipeline\n\n"
        "       Connection: JDBC + Key-Pair JWT Auth\n"
        "       Role: MENDIX_SERVICE_ROLE (51 grants, least-privilege)",
        top=Inches(1.5), font_size=Pt(13), color=LIGHT_GRAY)
    
    # ========== SLIDE 5: CoCo CLI Usage ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Built 100% with CoCo CLI")
    add_table(slide, 
        ["Task", "CoCo CLI Feature"],
        [
            ["16 stored procedures with retry logic", "SQL generation + execution"],
            ["Hybrid cross-check (rule + AI)", "Stored Proc design + testing"],
            ["228+ reference data records", "Bulk SQL generation"],
            ["Streamlit dashboard (5 pages)", "File management + deployment"],
            ["Marine weather integration", "Marketplace exploration + JOINs"],
            ["Snowpark analytics pipeline", "Python code generation"],
            ["SAP simulation (4 tables)", "Table DDL + procedure creation"],
            ["Cortex Agent (multilingual)", "Agent spec + configuration"],
        ],
        col_widths=[Inches(4.5), Inches(4.0)])
    
    # ========== SLIDE 6: Classification ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Feature 1: Auto-Classification")
    add_text(slide, "CALL CLASSIFY_DOCUMENT_TEXT('BILL OF LADING No. MAEU123...');\n> {document_type: \"BILL_OF_LADING\", confidence: 0.95}", top=Inches(1.5), font_size=Pt(12), color=ACCENT_BLUE)
    add_bullet(slide, [
        "17 maritime document types supported",
        "Confidence scoring - human-in-the-loop for < 85%",
        "Retry mechanism with exponential backoff (1s, 2s, 4s)",
        "All calls logged to AI_CALL_LOG for audit + cost tracking",
        "Model: llama3-8b (ultra-low cost)",
    ], top=Inches(3.0))
    
    # ========== SLIDE 7: Cross-Check ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Feature 2: Hybrid Cross-Check")
    add_text(slide, "8 SQL Rules (FREE)                     AI Fuzzy Match (edge cases only)", top=Inches(1.5), font_size=Pt(14), color=ACCENT_GREEN)
    add_bullet(slide, [
        "Rule 1: Weight discrepancy (>2% difference)",
        "Rule 2: Package count mismatch",
        "Rule 3: Vessel name mismatch",
        "Rule 4: Voyage number mismatch",
        "Rule 5: ETD date mismatch",
        "Rule 6: Incoterms mismatch",
        "Rule 7: Volume/CBM discrepancy",
        "Rule 8: Container number mismatch",
        "",
        "AI only for party name: \"VN SEAFOOD JSC\" = \"VIETNAM SEAFOOD JSC\" (saves 90% tokens)",
    ], top=Inches(2.3), font_size=Pt(13))
    
    # ========== SLIDE 8: Compliance ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Feature 3: Compliance Engine")
    add_text(slide, "CALL CHECK_COMPLIANCE(1);  -- 100% Deterministic, zero AI hallucination", top=Inches(1.5), font_size=Pt(12), color=ACCENT_BLUE)
    add_table(slide, 
        ["Check", "Reference Data", "Result"],
        [
            ["HS Code valid?", "138 codes (97 chapters + 4-digit)", "PASS/FAIL"],
            ["Dangerous Goods?", "DG flag on HS_CODE_REFERENCE", "WARNING"],
            ["VGM present?", "SOLAS requirement", "PASS/FAIL"],
            ["Route-specific docs?", "EU->CoO, US->ISF, JP->NACCS", "WARNING"],
        ],
        top=Inches(2.5),
        col_widths=[Inches(2.5), Inches(3.5), Inches(2.0)])
    
    # ========== SLIDE 9: Fraud Detection ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Feature 4: Fraud Detection")
    add_text(slide, "5-Rule Anomaly Detection (Pure SQL, Instant, Zero AI Cost)", top=Inches(1.4), font_size=Pt(14), color=ACCENT_YELLOW)
    add_table(slide,
        ["Rule", "Description", "Severity"],
        [
            ["DUPLICATE_BL", "Same B/L number on multiple documents", "HIGH"],
            ["DUPLICATE_CONTAINER", "Same container on different B/Ls", "HIGH"],
            ["INVALID_CONTAINER", "Fails ISO 6346 check-digit", "MEDIUM"],
            ["WEIGHT_ANOMALY", "Weight/volume ratio abnormal", "MEDIUM"],
            ["POSSIBLE_COPY", "Same shipper+consignee+weight+date", "HIGH"],
        ],
        top=Inches(2.3),
        col_widths=[Inches(2.5), Inches(4.0), Inches(1.5)])
    
    # ========== SLIDE 10: Container Photo ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Feature 5: Container Photo Verification")
    add_text(slide, "CALL VERIFY_CONTAINER_PHOTO('@stage/photo.jpg', 'MAEU1234567');\n> {container_match: true, seal_match: true, condition: \"good\"}", top=Inches(1.5), font_size=Pt(12), color=ACCENT_BLUE)
    add_bullet(slide, [
        "AI_PARSE_DOCUMENT for OCR - extracts container number from photo",
        "Cross-references against B/L container number",
        "Seal number integrity verification",
        "Container physical condition assessment",
        "Used at gate entry - trucks verified before entering port",
    ], top=Inches(3.0))
    
    # ========== SLIDE 11: Marketplace ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Marketplace: 3 Live Data Sources (Zero ETL)")
    add_table(slide,
        ["Source", "Data", "Workflow Use"],
        [
            ["Pelmorex Weather", "3,324 port forecasts", "Delay alerts"],
            ["Public Free - FX", "12 currency pairs (live)", "SAP charge conversion"],
            ["Public Free - WTO", "VN + 7 partner trade stats", "Route demand"],
            ["Public Free - ITA", "1,816 restricted entities", "Sanction screening"],
        ],
        top=Inches(1.8),
        col_widths=[Inches(2.8), Inches(3.2), Inches(2.5)])
    add_text(slide, "All free | Auto-refresh | Zero infrastructure | Zero-copy sharing", top=Inches(4.5), font_size=Pt(12), color=ACCENT_BLUE)
    
    # ========== SLIDE 12: Streamlit ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Analytics: Streamlit-in-Snowflake")
    add_kpi_box(slide, "Overview", "KPIs", Inches(0.5), Inches(1.6), Inches(1.7))
    add_kpi_box(slide, "Documents", "Explorer", Inches(2.4), Inches(1.6), Inches(1.7))
    add_kpi_box(slide, "Compliance", "Monitor", Inches(4.3), Inches(1.6), Inches(1.7))
    add_kpi_box(slide, "Fraud", "Alerts", Inches(6.2), Inches(1.6), Inches(1.7))
    add_kpi_box(slide, "AI Analytics", "Cost", Inches(8.1), Inches(1.6), Inches(1.7))
    add_bullet(slide, [
        "5-page dashboard deployed natively in Snowflake",
        "No external hosting required",
        "Real-time data from operational tables",
        "AI cost monitoring - essential for enterprise budget control",
        "Object: MENDIX_APP.AGENTS.VF_LOGISTICS_DASHBOARD",
    ], top=Inches(3.3))
    
    # ========== SLIDE 13: Security ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Security & Governance")
    add_table(slide,
        ["Feature", "Implementation"],
        [
            ["Least-Privilege Role", "MENDIX_SERVICE_ROLE: 51 grants, SELECT-only"],
            ["Authentication", "Key-pair JWT (no password)"],
            ["Audit Trail", "AI_CALL_LOG: every AI call recorded"],
            ["Cost Control", "Token + latency tracking per procedure"],
            ["Agent Data Policy", "Specific lookup OK, bulk export BLOCKED"],
            ["Anti-Scraping", "Broad queries return summaries only"],
        ],
        top=Inches(1.8),
        col_widths=[Inches(3.0), Inches(5.5)])
    
    # ========== SLIDE 14: AI Agent ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "AI Agent: Professional Logistics Search")
    add_table(slide,
        ["Search Type", "Examples"],
        [
            ["By Identifier", "B/L number, Container ISO, Booking ref, IMO"],
            ["By Route", "Port pair (VNSGN>JPYOK), Country, Trade lane"],
            ["By Party", "Shipper, Consignee, Carrier name"],
            ["By Time", "ETD/ETA range, Date of issue"],
            ["By Cargo", "Commodity, HS Code, Container type, Weight"],
            ["By Status", "Document / Shipment / Compliance status"],
        ],
        top=Inches(1.8),
        col_widths=[Inches(2.5), Inches(6.0)])
    add_text(slide, "Domain: Incoterms 2020 | IMDG Code | VGM/SOLAS | Demurrage | L/C | Carrier Alliances | B/L Types", top=Inches(5.5), font_size=Pt(11), color=MID_GRAY)
    
    # ========== SLIDE 14: Business Impact ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Business Impact")
    add_table(slide,
        ["Metric", "Before", "After", "Improvement"],
        [
            ["Classification time", "15-30 min", "2-3 sec", "99.7% faster"],
            ["Compliance coverage", "60% manual", "100% auto", "Full coverage"],
            ["Cross-check accuracy", "Error-prone", "Rule+AI hybrid", "Zero missed"],
            ["Fraud detection", "Reactive", "Proactive", "Early detection"],
            ["AI cost per document", "-", "~$0.001", "Ultra-low cost"],
            ["SAP posting time", "Manual entry", "Automatic", "Zero manual work"],
        ],
        top=Inches(1.8),
        col_widths=[Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.0)])
    
    # ========== SLIDE 15: SAP Integration ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "SAP S/4HANA Integration (Phase 4)")
    add_text(slide, "B/L Approved > Automatic SAP Postings:", top=Inches(1.5), font_size=Pt(14), color=WHITE)
    add_bullet(slide, [
        "FI: Vendor Invoice (Debit Freight 4210000 / Credit AP 2100000)",
        "MM: Goods Receipt (MIGO 101 > Plant VF01, Storage WH01)",
        "SD: Delivery + Billing (customer invoice generated)",
        "CO: Cost Allocation (Ocean, THC, Doc Fee, BAF breakdown)",
        "",
        "Current: SAP tables simulated in Snowflake (4 tables + 4 procedures)",
        "Future: SAP No-Copy (Datasphere federation, zero ETL)",
    ], top=Inches(2.5), font_size=Pt(14))
    
    # ========== SLIDE 16: Demo & Roadmap ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Live Demo & Roadmap")
    add_text(slide, "Demo Flow (6 steps):", top=Inches(1.4), font_size=Pt(14), color=WHITE, bold=True)
    add_bullet(slide, [
        "1. Upload document > Auto-classify (B/L, 95% confidence)",
        "2. Cross-check B/L vs Invoice > Weight discrepancy found",
        "3. Compliance scan > DG flagged, VGM verified",
        "4. Fraud scan > Duplicate container detected",
        "5. Weather check > Strong wind at destination port",
        "6. SAP posting > FI + MM + SD + CO created automatically",
    ], top=Inches(2.2), font_size=Pt(14))
    add_table(slide,
        ["Phase", "Scope", "Status"],
        [
            ["Phase 1", "AI Document Intelligence", "LIVE"],
            ["Phase 2", "Gate Management (3,000 trucks/day)", "Designed"],
            ["Phase 3", "Warehouse & Yard (7 DCs)", "Designed"],
            ["Phase 4", "SAP S/4HANA (No-Copy)", "Simulated"],
            ["Phase 5", "Multi-Tenant Data Isolation", "Planned"],
        ],
        top=Inches(5.0),
        col_widths=[Inches(1.5), Inches(4.5), Inches(2.0)])
    
    # ========== SLIDE 17: Team ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "Team SORA", font_size=Pt(40))
    add_table(slide,
        ["Member", "Role", "Contact"],
        [
            ["Chau Phuoc Hoa", "Team Lead / Backend Developer", "hoachauphuoc@gmail.com"],
            ["Nguyen Quoc Cuong", "Frontend Developer", "walkeralan620@gmail.com"],
        ],
        top=Inches(2.0),
        col_widths=[Inches(3.0), Inches(3.5), Inches(3.0)])
    add_text(slide, "Project: VF Logistics AI-Powered Seaport Platform\nHackathon: Snowflake CoCo CLI Hackathon 2026\nTech Stack: Snowflake (Cortex AI, Snowpark, Streamlit) + Mendix", top=Inches(4.0), font_size=Pt(14), color=MID_GRAY)
    
    # Save
    output_path = r"C:\Users\phuochoa\Mendix\VF_Logistics_Portal-main\snowflake-backend\VF_Logistics_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
